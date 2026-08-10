import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme colors
    COLOR_BG = RGBColor(10, 14, 23)         # #0a0e17
    COLOR_SURFACE = RGBColor(17, 24, 39)     # #111827
    COLOR_PRIMARY = RGBColor(59, 130, 246)   # #3b82f6
    COLOR_TEXT = RGBColor(241, 245, 249)     # #f1f5f9
    COLOR_MUTED = RGBColor(148, 163, 184)    # #94a3b8
    COLOR_ACCENT = RGBColor(6, 182, 212)     # #06b6d4
    COLOR_GREEN = RGBColor(16, 185, 129)     # #10b981
    COLOR_AMBER = RGBColor(245, 158, 11)     # #f59e0b
    COLOR_RED = RGBColor(239, 68, 68)       # #ef4444

    def add_blank_slide(title_text, category_text="RENEWTWIN"):
        slide = prs.slides.add_slide(blank_layout)
        
        # Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # Header bar
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ACCENT
        p0.font.name = 'Arial'

        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(28)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_TEXT
        p1.font.name = 'Arial'

        return slide

    def add_card(slide, left, top, width, height, title, items, border_color=COLOR_PRIMARY):
        # Card background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_SURFACE
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(height - 0.4))
        tf = txBox.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(18)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_ACCENT
        p0.font.name = 'Arial'

        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_TEXT
            p.font.name = 'Arial'
            p.space_before = Pt(8)

    # ------------------- SLIDE 1: Title -------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()

    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "RENEWTWIN"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Digital Twin for Predictive Management of Renewable Energy Assets"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_TEXT
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = "Energy Innovation Challenge 2026 | Track 4: Digital Asset Management\nMC²Plus × Oil India Ltd. × IIT Kharagpur"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_MUTED
    p3.space_before = Pt(24)

    # ------------------- SLIDE 2: Problem -------------------
    slide2 = add_blank_slide("The Operational Problem in Renewable Plants", "PROBLEM STATEMENT")
    add_card(slide2, 0.8, 1.6, 5.6, 5.2, "Massive Asset Scale & Distribution", [
        "Modern solar plants span thousands of acres with 50,000+ PV modules.",
        "Manual physical inspections are slow, expensive, and logistically difficult.",
        "Defects occur silently due to thermal stress, micro-cracks, and environmental exposure."
    ], COLOR_RED)
    add_card(slide2, 6.8, 1.6, 5.6, 5.2, "Impact of Delayed Fault Detection", [
        "Avoidable energy losses occur long before manual discovery.",
        "Unaddressed hot-spots lead to permanent module failure & fire hazards.",
        "Maintenance is reactive rather than predictive, driving up O&M costs."
    ], COLOR_RED)

    # ------------------- SLIDE 3: Existing Limitations -------------------
    slide3 = add_blank_slide("Why Existing Solutions Fall Short", "CURRENT LIMITATIONS")
    add_card(slide3, 0.8, 1.6, 3.6, 5.2, "Traditional SCADA", [
        "Monitors aggregate plant output only.",
        "Lacks module-level granular insight.",
        "Cannot isolate root cause of degradation."
    ], COLOR_MUTED)
    add_card(slide3, 4.8, 1.6, 3.6, 5.2, "Periodic Drone Imagery", [
        "Conducted only 1-2 times per year.",
        "Generates static imagery without live telemetry linkage.",
        "Leaves long blind spots between inspections."
    ], COLOR_MUTED)
    add_card(slide3, 8.8, 1.6, 3.6, 5.2, "Isolated ML Tools", [
        "Perform single-image classification.",
        "No digital state persistence.",
        "Fails to prioritize tasks based on total energy impact."
    ], COLOR_MUTED)

    # ------------------- SLIDE 4: Solution Overview -------------------
    slide4 = add_blank_slide("RenewTwin: Continuous Intelligence via Digital Twins", "THE SOLUTION")
    add_card(slide4, 0.8, 1.6, 5.6, 5.2, "Core Conceptual Shift", [
        "TRADITIONAL: Periodic inspection -> Fault discovered -> Maintenance -> Energy loss already occurred.",
        "RENEWTWIN: Continuous telemetry + AI defect detection -> Digital Twin state persistence -> Asset Health Scoring -> Automated Priority Dispatch."
    ], COLOR_GREEN)
    add_card(slide4, 6.8, 1.6, 5.6, 5.2, "Key Platform Capabilities", [
        "Renewable Asset Digital Twin Engine",
        "Computer Vision Surface Defect Detection (ResNet-18)",
        "Operational Telemetry Anomaly Detection (Isolation Forest)",
        "Explainable Asset Health Scoring & Priority Ranking"
    ], COLOR_PRIMARY)

    # ------------------- SLIDE 5: System Architecture -------------------
    slide5 = add_blank_slide("End-to-End System Architecture", "SYSTEM ARCHITECTURE")
    add_card(slide5, 0.8, 1.6, 11.7, 5.2, "Multimodal Architecture Pipeline", [
        "Data Ingestion: High-resolution visual imagery + Operational telemetry (Power, Temp, Irradiance).",
        "AI Analysis Layer: ResNet-18 (Defect Classifier) + Isolation Forest (Operational Anomaly Detector).",
        "Digital Twin Engine: Persistent digital representation tracking state, history, and physical parameters.",
        "Health & Priority Engine: 4-factor Asset Health Index (AHI) computation.",
        "Operator Interface: FastAPI Backend + Real-time React Industrial Monitoring Dashboard."
    ], COLOR_ACCENT)

    # ------------------- SLIDE 6: AI/ML Methodology -------------------
    slide6 = add_blank_slide("Dual AI Engine: Vision + Telemetry", "AI & ML METHODOLOGY")
    add_card(slide6, 0.8, 1.6, 5.6, 5.2, "Visual Defect Detection (CNN)", [
        "Model: Deep Transfer Learning with ResNet-18.",
        "Target Classes: Clean (None), Micro-Crack, Hotspot, Inactive/Cell Damage.",
        "Outputs: Multi-class defect probability & confidence score."
    ], COLOR_PRIMARY)
    add_card(slide6, 6.8, 1.6, 5.6, 5.2, "Telemetry Anomaly Detection", [
        "Model: Isolation Forest (scikit-learn).",
        "Features: Power deviation ratio (actual/expected), Temperature elevation, Irradiance.",
        "Outputs: Unsupervised anomaly score (0.0 to 1.0)."
    ], COLOR_ACCENT)

    # ------------------- SLIDE 7: Digital Twin Design -------------------
    slide7 = add_blank_slide("Intelligent Digital Twin Schema", "DIGITAL TWIN DESIGN")
    add_card(slide7, 0.8, 1.6, 11.7, 5.2, "Digital Twin Object Schema", [
        "Identity & Location: asset_id (e.g. PV-A-014), asset_type, location (Array A/Row 17).",
        "Physical Metrics: rated_capacity_kw (400 kW), current_power_kw, expected_power_kw, temperature_c.",
        "AI Intelligence State: defect_class, defect_probability, anomaly_score.",
        "Health & Action: health_score (53.8), risk_level (AT RISK), maintenance_priority (#1), recommended_action."
    ], COLOR_GREEN)

    # ------------------- SLIDE 8: Health Score & Risk -------------------
    slide8 = add_blank_slide("Prototype Asset Health Index (AHI)", "HEALTH & RISK SCORING")
    add_card(slide8, 0.8, 1.6, 5.6, 5.2, "Scoring Formula Breakdown", [
        "AHI = 100 - Penalties",
        "Visual Defect Penalty: 40% Weight",
        "Operational Anomaly Penalty: 30% Weight",
        "Performance Deviation Penalty: 20% Weight",
        "Thermal Elevation Penalty: 10% Weight"
    ], COLOR_AMBER)
    add_card(slide8, 6.8, 1.6, 5.6, 5.2, "Operational Risk Categories", [
        "HEALTHY (90 - 100): Normal operation, routine monitoring.",
        "MONITOR (75 - 89): Minor deviation/soiling, schedule periodic check.",
        "AT RISK (50 - 74): Thermal anomaly / defect detected, prioritize inspection.",
        "CRITICAL (< 50): Severe failure imminent / heavy loss, immediate dispatch."
    ], COLOR_AMBER)

    # ------------------- SLIDE 9: Maintenance Prioritization -------------------
    slide9 = add_blank_slide("Actionable Maintenance Prioritization", "MAINTENANCE ENGINE")
    add_card(slide9, 0.8, 1.6, 11.7, 5.2, "Impact-Driven Task Ranking", [
        "Answers: 'Which specific panel requires field technician intervention right now?'",
        "Consolidates failure probability, estimated energy loss percentage, and thermal risk.",
        "Outputs dynamically ordered dispatch queue (#1 PV-A-014: Hotspot, 46.9% loss estimated).",
        "Eliminates wasted technician trips and focuses labor on maximum ROI repairs."
    ], COLOR_PRIMARY)

    # ------------------- SLIDE 10: Dashboard Demo -------------------
    slide10 = add_blank_slide("Industrial Operator Dashboard MVP", "PRODUCT DEMONSTRATION")
    add_card(slide10, 0.8, 1.6, 5.6, 5.2, "Full-Stack Implementation", [
        "Backend: Python FastAPI with SQLite aiosqlite database & RESTful routes.",
        "Frontend: React + Vite with dark-themed industrial aesthetic.",
        "Live Telemetry & Digital Twin inspector modal with real-time update."
    ], COLOR_PRIMARY)
    add_card(slide10, 6.8, 1.6, 5.6, 5.2, "Demo Flow Highlights", [
        "1. Plant KPI Overview (24 Active PV Assets).",
        "2. Risk Distribution breakdown (13 Healthy, 5 Monitor, 6 At Risk).",
        "3. Interactive Digital Twin drill-down on critical assets.",
        "4. AI Inspection Upload simulation with live model feedback."
    ], COLOR_GREEN)

    # ------------------- SLIDE 11: Scalability & Future Scope -------------------
    slide11 = add_blank_slide("Future Roadmap & Technical Scalability", "SCALABILITY & ROADMAP")
    add_card(slide11, 0.8, 1.6, 5.6, 5.2, "Multi-Asset Extension", [
        "Expand digital twin schemas from Solar PV to Wind Turbines and Battery Storage (BESS).",
        "Edge Deployment: Lightweight ONNX/TensorRT inference on drone-mounted compute."
    ], COLOR_ACCENT)
    add_card(slide11, 6.8, 1.6, 5.6, 5.2, "Industrial Enterprise Integration", [
        "Time-Series DB: Migration from SQLite to InfluxDB / TimescaleDB.",
        "Enterprise Connectors: Automated work order generation for SAP / IBM Maximo EAM."
    ], COLOR_ACCENT)

    # ------------------- SLIDE 12: Conclusion -------------------
    slide12 = add_blank_slide("RenewTwin: Transforming Renewable O&M", "CONCLUSION & IMPACT")
    add_card(slide12, 0.8, 1.6, 11.7, 5.2, "Key Takeaways & Value Delivered", [
        "Predictive Management: From reactive repair to continuous AI asset health monitoring.",
        "Demonstrable MVP: Complete working backend, ML pipeline, digital twin engine, and React UI.",
        "Commercial Viability: Direct reduction in LCOE, O&M expenses, and avoidable power loss.",
        "GitHub Repository: https://github.com/bhaktictak/RenewTwin"
    ], COLOR_GREEN)

    output_path = "RenewTwin_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at {output_path}")

if __name__ == '__main__':
    create_presentation()
